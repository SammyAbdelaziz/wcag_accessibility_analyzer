"""
Object Model Normalization Layer
Wraps python-docx and python-pptx to provide higher-level document representations.
Enables richer semantic analysis and more accurate findings.

Use Cases:
  - Validate list coherence (nesting depth, bullet styles)
  - Extract heading hierarchy with semantic meaning
  - Detect table semantic completeness (header scope)
  - Analyze reading order vs. visual order
"""

from __future__ import annotations
from dataclasses import dataclass, field
from typing import List, Optional, Dict, Set, Tuple
from enum import Enum
import io


# ────────────────────────────────────────────────────────────────────
# DOCX Object Model
# ────────────────────────────────────────────────────────────────────

class ListStyle(str, Enum):
    BULLET = "bullet"
    NUMBERED = "numbered"
    MULTILEVEL = "multilevel"
    UNKNOWN = "unknown"


@dataclass
class DocxHeadingInfo:
    """Heading extracted from DOCX with semantic context."""
    level: int  # 1-6
    text: str
    paragraph_index: int
    has_style: bool  # True if Heading 1/2/3, etc.; False if styled-as-heading
    font_size_pt: Optional[int]
    is_bold: bool


@dataclass
class DocxListInfo:
    """List structure with coherence validation."""
    list_id: str  # Unique identifier in document
    style: ListStyle
    items: List[str]  # Text of each list item
    nesting_levels: List[int]  # Level (0, 1, 2, ...) of each item
    paragraph_indices: List[int]  # Paragraph indices
    is_coherent: bool = True  # True if nesting is valid (no gaps, consistent)
    max_depth: int = 0
    violations: List[str] = field(default_factory=list)  # "skip at item 3", etc.
    # Reserved for future detection of formatting-only lists.
    # Populated only when a detector identifies non-semantic list patterns.
    formatting_type: Optional[str] = None
    accessibility_issues: List[Dict] = field(default_factory=list)
    
    def validate_coherence(self) -> bool:
        """Check if list nesting is coherent (no skips, logical progression)."""
        if not self.nesting_levels:
            return True
        
        seen_levels = set()
        for i, level in enumerate(self.nesting_levels):
            if level > max(seen_levels or {0}) + 1:
                # Skipped a level
                self.violations.append(f"Skipped level {max(seen_levels or {0}) + 1} at item {i+1}")
                self.is_coherent = False
            seen_levels.add(level)
            self.max_depth = max(self.max_depth, level)
        
        return self.is_coherent


@dataclass
class DocxTableInfo:
    """Table with semantic completeness validation."""
    index: int
    row_count: int
    col_count: int
    has_header_row: bool
    header_text: List[str]  # Text of header cells (if present)
    paragraph_index: int
    has_merged_cells: bool = False
    first_col_is_header: bool = False  # True if first column appears to be row headers
    is_complete: bool = True  # All cells have content
    
    def validate_semantics(self) -> Dict[str, bool]:
        """Check semantic completeness."""
        return {
            'has_header': self.has_header_row,
            'has_content': self.is_complete,
            'has_merged_cells': self.has_merged_cells,
            'likely_row_headers': self.first_col_is_header,
        }


@dataclass
class DocxNormalizedModel:
    """High-level DOCX document representation."""
    filename: str
    title: Optional[str]
    language: Optional[str]
    headings: List[DocxHeadingInfo]
    lists: List[DocxListInfo]
    tables: List[DocxTableInfo]
    paragraph_count: int
    image_count: int
    hyperlink_count: int
    # Future-facing collection: lists detected as visually-formatted but not
    # semantically marked. Defaults to empty until a detector populates it.
    formatted_lists: List[DocxListInfo] = field(default_factory=list)
    
    # Derived data
    heading_hierarchy_valid: bool = True
    all_lists_coherent: bool = True
    heading_levels_present: Set[int] = field(default_factory=set)
    
    def validate_structure(self) -> Dict[str, any]:
        """Validate overall document structure."""
        self.heading_levels_present = {h.level for h in self.headings}
        
        # Check heading hierarchy
        seen = set()
        for level in sorted(self.heading_levels_present):
            if level > max(seen or {0}) + 1:
                self.heading_hierarchy_valid = False
            seen.add(level)
        
        # Check list coherence
        self.all_lists_coherent = all(l.is_coherent for l in self.lists)
        
        return {
            'heading_hierarchy_valid': self.heading_hierarchy_valid,
            'all_lists_coherent': self.all_lists_coherent,
            'heading_levels': sorted(self.heading_levels_present),
            'list_count': len(self.lists),
            'table_count': len(self.tables),
        }


class DocxNormalizer:
    """Extract normalized object model from DOCX using python-docx."""
    
    def __init__(self, docx_bytes: bytes, filename: str):
        self.docx_bytes = docx_bytes
        self.filename = filename
        self.doc = None
        self._load_document()
    
    def _load_document(self):
        """Load DOCX using python-docx."""
        try:
            from docx import Document
            self.doc = Document(io.BytesIO(self.docx_bytes))
        except Exception:
            # python-docx may fail on minimal/non-standard packages; rules
            # using this normalizer treat None as "skip gracefully".
            self.doc = None
    
    def normalize(self) -> Optional[DocxNormalizedModel]:
        """Extract normalized model from document."""
        if not self.doc:
            return None
        
        title = self._extract_title()
        language = self._extract_language()
        headings = self._extract_headings()
        lists = self._extract_lists()
        tables = self._extract_tables()
        
        image_count = sum(1 for rel in self.doc.part.rels.values() 
                         if 'image' in rel.target_ref)
        hyperlink_count = sum(1 for para in self.doc.paragraphs 
                             for _ in para.hyperlinks)
        
        model = DocxNormalizedModel(
            filename=self.filename,
            title=title,
            language=language,
            headings=headings,
            lists=lists,
            tables=tables,
            paragraph_count=len(self.doc.paragraphs),
            image_count=image_count,
            hyperlink_count=hyperlink_count,
        )
        
        model.validate_structure()
        return model
    
    def _extract_title(self) -> Optional[str]:
        """Extract document title from core properties."""
        try:
            return self.doc.core_properties.title or None
        except Exception:
            return None
    
    def _extract_language(self) -> Optional[str]:
        """Extract document language."""
        try:
            # Language is typically in element.xml settings, not exposed by python-docx
            # This is a placeholder; actual implementation would parse XML
            return None
        except Exception:
            return None
    
    def _extract_headings(self) -> List[DocxHeadingInfo]:
        """Extract all headings with semantic context."""
        headings = []
        para_index = 0
        
        for paragraph in self.doc.paragraphs:
            style = paragraph.style
            style_name = style.name if style else "Normal"
            
            # Check if this is a Heading style
            if style_name.startswith('Heading'):
                try:
                    level = int(style_name.split()[-1])
                except (ValueError, IndexError):
                    level = 0
                
                font_size = None
                is_bold = False
                
                # Extract run properties from first run
                for run in paragraph.runs:
                    if run.font.size:
                        font_size = run.font.size.pt
                    if run.bold:
                        is_bold = True
                    break
                
                headings.append(DocxHeadingInfo(
                    level=level,
                    text=paragraph.text.strip(),
                    paragraph_index=para_index,
                    has_style=True,
                    font_size_pt=font_size,
                    is_bold=is_bold,
                ))
            
            para_index += 1
        
        return headings
    
    def _extract_lists(self) -> List[DocxListInfo]:
        """Extract list structures with coherence validation."""
        lists = {}  # list_id -> DocxListInfo
        para_index = 0
        
        for paragraph in self.doc.paragraphs:
            if paragraph.style and paragraph.style.name and paragraph.style.name.startswith('List'):
                # Determine list ID and level
                list_id = str(id(paragraph._element.getparent()))  # Simplified
                level = paragraph.paragraph_format.left_indent or 0
                
                if list_id not in lists:
                    lists[list_id] = DocxListInfo(
                        list_id=list_id,
                        style=ListStyle.BULLET,
                        items=[],
                        nesting_levels=[],
                        paragraph_indices=[],
                    )
                
                lists[list_id].items.append(paragraph.text.strip())
                lists[list_id].nesting_levels.append(int(level / 720))  # EMU to level
                lists[list_id].paragraph_indices.append(para_index)
            
            para_index += 1
        
        # Validate coherence
        for list_info in lists.values():
            list_info.validate_coherence()
        
        return list(lists.values())
    
    def _extract_tables(self) -> List[DocxTableInfo]:
        """Extract table information with semantic context."""
        tables = []
        para_index = 0
        table_index = 0
        
        for element in self.doc.element.body:
            tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag
            
            if tag == 'p':
                para_index += 1
            elif tag == 'tbl':
                try:
                    from docx.table import Table
                    table_element = Table(element, self.doc)
                    rows = table_element.rows
                    cols = len(rows[0].cells) if rows else 0
                    
                    # Extract header if present
                    header_text = []
                    has_header = False
                    if rows and hasattr(rows[0], '_tr'):
                        # Try to detect if first row is header
                        first_row = rows[0]
                        header_text = [cell.text.strip() for cell in first_row.cells]
                        has_header = bool(header_text[0])  # Simplified
                    
                    tables.append(DocxTableInfo(
                        index=table_index,
                        row_count=len(rows),
                        col_count=cols,
                        has_header_row=has_header,
                        header_text=header_text,
                        paragraph_index=para_index,
                    ))
                    table_index += 1
                except Exception as e:
                    print(f"Warning: Failed to extract table {table_index}: {e}")
        
        return tables


# ────────────────────────────────────────────────────────────────────
# PPTX Object Model
# ────────────────────────────────────────────────────────────────────

@dataclass
class PptxShapeInfo:
    """Shape with semantic role information."""
    shape_id: int
    shape_name: str
    shape_type: str  # "picture", "text", "title", "placeholder", "group"
    placeholder_role: Optional[str]  # "title", "body", "ctrTitle", etc.
    has_alt_text: bool
    alt_text: Optional[str]
    text_content: Optional[str]
    z_order: int
    is_decorative: bool = False


@dataclass
class PptxSlideInfo:
    """Slide with master/layout context."""
    slide_index: int
    slide_number: int
    title: Optional[str]
    has_title_placeholder: bool
    shapes: List[PptxShapeInfo]
    is_title_slide: bool = False
    has_content: bool = True


@dataclass
class PptxNormalizedModel:
    """High-level PPTX presentation representation."""
    filename: str
    title: Optional[str]
    slide_count: int
    slides: List[PptxSlideInfo]
    all_slides_have_titles: bool = True
    
    def validate_structure(self) -> Dict[str, any]:
        """Validate presentation structure."""
        self.all_slides_have_titles = all(
            slide.has_title_placeholder 
            for i, slide in enumerate(self.slides)
            if i > 0  # Skip title slide
        )
        
        return {
            'all_slides_have_titles': self.all_slides_have_titles,
            'slide_count': self.slide_count,
            'shapes_total': sum(len(s.shapes) for s in self.slides),
        }


class PptxNormalizer:
    """Extract normalized object model from PPTX using python-pptx."""
    
    def __init__(self, pptx_bytes: bytes, filename: str):
        self.pptx_bytes = pptx_bytes
        self.filename = filename
        self.prs = None
        self._load_presentation()
    
    def _load_presentation(self):
        """Load PPTX using python-pptx."""
        try:
            from pptx import Presentation
            self.prs = Presentation(io.BytesIO(self.pptx_bytes))
        except Exception as e:
            print(f"Warning: Failed to load PPTX with python-pptx: {e}")
            self.prs = None
    
    def normalize(self) -> Optional[PptxNormalizedModel]:
        """Extract normalized model from presentation."""
        if not self.prs:
            return None
        
        title = self._extract_title()
        slides = self._extract_slides()
        
        model = PptxNormalizedModel(
            filename=self.filename,
            title=title,
            slide_count=len(self.prs.slides),
            slides=slides,
        )
        
        model.validate_structure()
        return model
    
    def _extract_title(self) -> Optional[str]:
        """Extract presentation title."""
        try:
            return self.prs.core_properties.title or None
        except Exception:
            return None
    
    def _extract_slides(self) -> List[PptxSlideInfo]:
        """Extract slide information with master/layout context."""
        slides = []
        
        for slide_idx, slide in enumerate(self.prs.slides):
            title = None
            has_title_placeholder = False
            shapes = []
            
            for z_order, shape in enumerate(slide.shapes):
                shape_info = PptxShapeInfo(
                    shape_id=shape.shape_id,
                    shape_name=shape.name,
                    shape_type=self._get_shape_type(shape),
                    placeholder_role=self._get_placeholder_role(shape),
                    has_alt_text=bool(getattr(shape, 'alt_text', None)),
                    alt_text=getattr(shape, 'alt_text', None),
                    text_content=shape.text if hasattr(shape, 'text') else None,
                    z_order=z_order,
                )
                
                shapes.append(shape_info)
                
                # Check if this is the title
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    phf = shape.placeholder_format
                    if phf.type == 1:  # PP_PLACEHOLDER.TITLE
                        title = shape.text if hasattr(shape, 'text') else None
                        has_title_placeholder = True
            
            is_title_slide = slide_idx == 0
            slide_info = PptxSlideInfo(
                slide_index=slide_idx,
                slide_number=slide_idx + 1,
                title=title,
                has_title_placeholder=has_title_placeholder,
                shapes=shapes,
                is_title_slide=is_title_slide,
            )
            slides.append(slide_info)
        
        return slides
    
    def _get_shape_type(self, shape) -> str:
        """Determine shape type."""
        if hasattr(shape, 'shape_type'):
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            st = shape.shape_type
            if st == MSO_SHAPE_TYPE.PICTURE:
                return "picture"
            elif st == MSO_SHAPE_TYPE.GROUP:
                return "group"
            elif st == MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT:
                return "ole"
        
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            return "placeholder"
        
        if hasattr(shape, 'text'):
            return "text"
        
        return "unknown"
    
    def _get_placeholder_role(self, shape) -> Optional[str]:
        """Get placeholder role."""
        try:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                phf = shape.placeholder_format
                if hasattr(phf, 'type'):
                    # Map type to role name
                    from pptx.enum.shapes import PP_PLACEHOLDER
                    return {
                        PP_PLACEHOLDER.TITLE: "title",
                        PP_PLACEHOLDER.BODY: "body",
                        PP_PLACEHOLDER.CENTER_TITLE: "ctrTitle",
                        PP_PLACEHOLDER.SUBTITLE: "subtitle",
                    }.get(phf.type, "unknown")
        except Exception:
            pass
        
        return None
