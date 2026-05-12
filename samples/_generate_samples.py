"""Generate tiny, non-sensitive sample docs for the POC smoke test."""
from pathlib import Path
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PInches, Pt

OUT = Path(__file__).parent

# --- sample.docx ---------------------------------------------------------
d = Document()
d.add_heading("Sample Accessibility Test Document", level=1)
d.add_paragraph(
    "This is a tiny synthetic document used to smoke-test the WCAG analyzer. "
    "It contains a heading, a paragraph, a list, and a small table so the "
    "analyzer has something to report on."
)
d.add_paragraph("Key points:", style="Intense Quote")
for item in ("Headings present", "Plain text paragraph", "Bullet list", "Small data table"):
    d.add_paragraph(item, style="List Bullet")
table = d.add_table(rows=2, cols=2)
table.style = "Light Grid Accent 1"
table.rows[0].cells[0].text = "Category"
table.rows[0].cells[1].text = "Value"
table.rows[1].cells[0].text = "Example"
table.rows[1].cells[1].text = "42"
d.save(OUT / "sample.docx")

# --- sample.pptx ---------------------------------------------------------
p = Presentation()
slide = p.slides.add_slide(p.slide_layouts[0])
slide.shapes.title.text = "Sample Accessibility Test Deck"
slide.placeholders[1].text = "Synthetic content for WCAG analyzer smoke testing"

slide = p.slides.add_slide(p.slide_layouts[1])
slide.shapes.title.text = "Content slide"
body = slide.placeholders[1].text_frame
body.text = "First bullet point"
body.add_paragraph().text = "Second bullet point"
body.add_paragraph().text = "Third bullet point"
p.save(OUT / "sample.pptx")

print("Created sample.docx and sample.pptx in", OUT)
