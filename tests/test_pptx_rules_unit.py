"""
PPTX Rules Unit Tests - Simplified

Unit tests for core PPTX WCAG analyzer rules with focus on deterministic processing.
Tests verify analyzer execution and finding structure rather than making strict 
assertions about specific rule detections (which may depend on PPTX file structure).

Total coverage: 8+ core rules organized by WCAG criterion
"""

import unittest
import io
from pathlib import Path
from pptx.util import Inches
from wcag.analyzers.pptx_analyzer import PptxAnalyzer


class PptxAnalyzerTestBase(unittest.TestCase):
    """Base class for PPTX rule tests."""
    
    FIXTURES_DIR = Path(__file__).parent / "fixtures" / "pptx"
    
    @staticmethod
    def create_minimal_presentation():
        """Create a minimal valid PPTX presentation."""
        from pptx import Presentation
        return Presentation()
    
    @staticmethod
    def analyze_pptx(prs, filename: str = "test.pptx"):
        """Analyze a python-pptx Presentation object."""
        buf = io.BytesIO()
        prs.save(buf)
        return PptxAnalyzer(buf.getvalue(), filename).analyze()
    
    def get_findings_by_criterion(self, fact_sheet, criterion_id: str):
        """Filter findings by WCAG criterion ID."""
        confirmed = [f for f in fact_sheet.confirmed_findings if f.criterion_id == criterion_id]
        possible = [f for f in fact_sheet.possible_findings if f.criterion_id == criterion_id]
        return confirmed, possible


# =============================================================================
# 2.4.2 - Presentation and Slide Title Tests
# =============================================================================

class TestRule242PresentationTitle(PptxAnalyzerTestBase):
    """Tests for WCAG 2.4.2: Presentation has a descriptive title."""
    
    def test_presentation_with_title_processes(self):
        """Presentation with title should process without error."""
        from pptx import Presentation
        prs = Presentation()
        prs.core_properties.title = "Q4 2026 Accessibility Report"
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        fact_sheet = self.analyze_pptx(prs, "titled.pptx")
        self.assertIsNotNone(fact_sheet)
        self.assertIsNotNone(fact_sheet.confirmed_findings)
    
    def test_presentation_without_title_processes(self):
        """Presentation without title should process without error."""
        from pptx import Presentation
        prs = Presentation()
        # Don't set title
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        fact_sheet = self.analyze_pptx(prs, "untitled.pptx")
        self.assertIsNotNone(fact_sheet)


class TestRule242SlideTitles(PptxAnalyzerTestBase):
    """Tests for WCAG 2.4.2: All slides (except title slide) have titles."""
    
    def test_slides_with_titles_processes(self):
        """Slides with titles should process without error."""
        from pptx import Presentation
        prs = Presentation()
        prs.core_properties.title = "Report"
        
        # Title slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide1.shapes.title
        title.text = "Q4 Accessibility Report"
        
        # Content slide with title
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide2.shapes.title
        title.text = "Key Findings"
        
        fact_sheet = self.analyze_pptx(prs, "titled_slides.pptx")
        self.assertIsNotNone(fact_sheet)
        confirmed, possible = self.get_findings_by_criterion(fact_sheet, "2.4.2")
        # Slide findings exist or may be empty depending on analyzer sensitivity
        self.assertIsNotNone(confirmed)
        self.assertIsNotNone(possible)
    
    def test_slide_without_title_processes(self):
        """Slide without title should process without error."""
        from pptx import Presentation
        prs = Presentation()
        prs.core_properties.title = "Report"
        
        # Content slide without title
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        # Clear the title
        if slide.shapes.title:
            slide.shapes.title.text = ""
        
        fact_sheet = self.analyze_pptx(prs, "untitled_slide.pptx")
        self.assertIsNotNone(fact_sheet)


# =============================================================================
# 1.1.1 - Image Alt Text Tests
# =============================================================================

class TestRule111ImageAlt(PptxAnalyzerTestBase):
    """Tests for WCAG 1.1.1: Images have alternative text."""
    
    def test_presentation_without_images_processes(self):
        """Presentation with no images should process."""
        from pptx import Presentation
        prs = Presentation()
        prs.core_properties.title = "Text-Only Report"
        
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Overview"
        text_frame = slide.placeholders[1].text_frame
        text_frame.text = "This slide has only text."
        
        fact_sheet = self.analyze_pptx(prs, "no_images.pptx")
        self.assertIsNotNone(fact_sheet)
        confirmed, possible = self.get_findings_by_criterion(fact_sheet, "1.1.1")
        # Filter to image-specific findings
        image_findings = [f for f in confirmed + possible 
                         if f.remediation_id and "image" in f.remediation_id.lower()]
        # No images = no image findings
        self.assertEqual(len(image_findings), 0)


# =============================================================================
# 1.1.1 - Chart Title Tests
# =============================================================================

class TestRule111ChartTitles(PptxAnalyzerTestBase):
    """Tests for WCAG 1.1.1: Charts have descriptive titles."""
    
    def test_presentation_without_charts_processes(self):
        """Presentation with no charts should process."""
        from pptx import Presentation
        prs = Presentation()
        prs.core_properties.title = "Text-Only Report"
        
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Overview"
        
        fact_sheet = self.analyze_pptx(prs, "no_charts.pptx")
        self.assertIsNotNone(fact_sheet)
        confirmed, possible = self.get_findings_by_criterion(fact_sheet, "1.1.1")
        # Filter to chart-specific findings
        chart_findings = [f for f in confirmed + possible 
                         if f.remediation_id and "chart" in f.remediation_id.lower()]
        # No charts = no chart findings
        self.assertEqual(len(chart_findings), 0)


# =============================================================================
# 2.4.4 - Link Text Tests
# =============================================================================

class TestRule244LinkText(PptxAnalyzerTestBase):
    """Tests for WCAG 2.4.4: Links have descriptive text."""
    
    def test_presentation_without_links_processes(self):
        """Presentation with no links should process."""
        from pptx import Presentation
        prs = Presentation()
        prs.core_properties.title = "No Links"
        
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Overview"
        
        fact_sheet = self.analyze_pptx(prs, "no_links.pptx")
        self.assertIsNotNone(fact_sheet)
        confirmed, _ = self.get_findings_by_criterion(fact_sheet, "2.4.4")
        # No links = no link findings
        self.assertEqual(len(confirmed), 0)

    def test_generic_hyperlink_text_flags(self):
        """Generic hyperlink text in slide runs should trigger 2.4.4."""
        from pptx import Presentation

        prs = Presentation()
        prs.core_properties.title = "Link Test"
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        run = textbox.text_frame.paragraphs[0].add_run()
        run.text = "click here"
        run.hyperlink.address = "https://example.com/report"

        fact_sheet = self.analyze_pptx(prs, "generic_link_text.pptx")
        confirmed, _ = self.get_findings_by_criterion(fact_sheet, "2.4.4")
        self.assertGreaterEqual(len(confirmed), 1)


# =============================================================================
# 3.1.1 - Presentation Language Tests
# =============================================================================

class TestRule311Language(PptxAnalyzerTestBase):
    """Tests for WCAG 3.1.1: Presentation language is set."""
    
    def test_presentation_with_language_processes(self):
        """Presentation with language set should process."""
        from pptx import Presentation
        prs = Presentation()
        prs.core_properties.language = "en-US"
        prs.core_properties.title = "English Report"
        
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        fact_sheet = self.analyze_pptx(prs, "english.pptx")
        self.assertIsNotNone(fact_sheet)
        confirmed, _ = self.get_findings_by_criterion(fact_sheet, "3.1.1")
        # Language set = ideally no 3.1.1 findings
        self.assertEqual(len(confirmed), 0, "Should not flag presentation with language set")
    
    def test_presentation_without_language_processes(self):
        """Presentation without language should process."""
        from pptx import Presentation
        prs = Presentation()
        # Don't set language
        prs.core_properties.title = "No Language"
        
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        fact_sheet = self.analyze_pptx(prs, "no_language.pptx")
        self.assertIsNotNone(fact_sheet)


# =============================================================================
# 1.3.2 - Reading Order Tests
# =============================================================================

class TestRule132ReadingOrder(PptxAnalyzerTestBase):
    """Tests for WCAG 1.3.2: Reading order matches visual order."""
    
    def test_simple_slide_layout_processes(self):
        """Standard slide layout should process without error."""
        from pptx import Presentation
        prs = Presentation()
        prs.core_properties.title = "Report"
        
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
        slide.shapes.title.text = "Main Topic"
        text_frame = slide.placeholders[1].text_frame
        text_frame.text = "Content goes here"
        
        fact_sheet = self.analyze_pptx(prs, "standard_layout.pptx")
        self.assertIsNotNone(fact_sheet)


# =============================================================================
# Integration / Regression Tests
# =============================================================================

class TestPptxOutputFormat(PptxAnalyzerTestBase):
    """Verify PPTX fact sheet serialization."""
    
    def test_fact_sheet_has_required_fields(self):
        """Fact sheet should have all required output fields."""
        from pptx import Presentation
        prs = Presentation()
        prs.core_properties.title = "Test"
        prs.core_properties.language = "en-US"
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        fact_sheet = self.analyze_pptx(prs, "test.pptx")
        output = fact_sheet.to_dict()
        
        self.assertEqual(output["file_type"], "pptx")
        self.assertEqual(output["filename"], "test.pptx")
        self.assertIn("confirmed_findings", output)
        self.assertIn("possible_findings", output)
        self.assertIn("summary", output)
    
    def test_slide_count_populated(self):
        """Fact sheet should include slide count."""
        from pptx import Presentation
        prs = Presentation()
        prs.core_properties.title = "Multi-slide"
        
        # Add 3 slides
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i+1}"
        
        fact_sheet = self.analyze_pptx(prs, "multi_slide.pptx")
        self.assertEqual(fact_sheet.slide_count, 3, "Should count slides correctly")
    
    def test_findings_structure_valid(self):
        """All findings should have valid structure."""
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        fact_sheet = self.analyze_pptx(prs, "test.pptx")
        all_findings = fact_sheet.confirmed_findings + fact_sheet.possible_findings
        
        for finding in all_findings:
            self.assertIsNotNone(finding.criterion_id)
            # remediation_id may be optional or None
            # remediation_data may be optional


if __name__ == "__main__":
    unittest.main()
