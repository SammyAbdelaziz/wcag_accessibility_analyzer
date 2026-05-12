"""Phase E — WCAG 4.1.2 Name/Role/Value for form/actionable controls.

DOCX: Existing 1.3.1 form-control-labels rule now also emits a 4.1.2
finding for content controls without w:alias.
PPTX: New rule flags hyperlinked shapes that have no descr/title attr
and no visible text inside the shape.
"""
import io
import unittest
import zipfile

from wcag.analyzers.docx_analyzer import DocxAnalyzer
from wcag.analyzers.pptx_analyzer import PptxAnalyzer


def _confirmed(fs, criterion: str):
    return [f for f in fs.confirmed_findings if f.criterion_id == criterion]


# ─────────────────────────────────────────────────────────────────────────────
# DOCX 4.1.2 — content control without w:alias
# ─────────────────────────────────────────────────────────────────────────────
def _build_docx_with_sdt(sdtPr_inner: str) -> bytes:
    """Build a minimal DOCX with one w:sdt whose w:sdtPr inner is provided."""
    doc_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
<w:body>
    <w:sdt>
        <w:sdtPr>
            {sdtPr_inner}
        </w:sdtPr>
        <w:sdtContent>
            <w:p><w:r><w:t>Sample</w:t></w:r></w:p>
        </w:sdtContent>
    </w:sdt>
</w:body>
</w:document>
"""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
            '</Types>',
        )
        zf.writestr(
            "_rels/.rels",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
            '</Relationships>',
        )
        zf.writestr("word/document.xml", doc_xml)
    return buf.getvalue()


class TestDocx412FormControlNames(unittest.TestCase):
    def test_fail_unlabeled_text_control(self):
        # text content control with no alias and no tag → 4.1.2 finding
        data = _build_docx_with_sdt("<w:text/>")
        fs = DocxAnalyzer(data, "fail.docx").analyze()
        findings = _confirmed(fs, "4.1.2")
        self.assertEqual(len(findings), 1)
        self.assertIn("accessible name", findings[0].issue.lower())
        self.assertEqual(findings[0].remediation_id, "form_control_names_412")

    def test_pass_labeled_text_control(self):
        # alias present → no 4.1.2 finding
        data = _build_docx_with_sdt(
            '<w:alias w:val="Patient name"/><w:text/>'
        )
        fs = DocxAnalyzer(data, "pass.docx").analyze()
        self.assertEqual(len(_confirmed(fs, "4.1.2")), 0)


# ─────────────────────────────────────────────────────────────────────────────
# PPTX 4.1.2 — hyperlinked shape without accessible name
# ─────────────────────────────────────────────────────────────────────────────
def _build_pptx_with_hyperlinked_shape(
    descr: str = "",
    has_text: bool = False,
) -> bytes:
    """Build a minimal pptx with a single rectangle shape that has an
    a:hlinkClick (hyperlink). Caller controls whether descr is set and
    whether the shape has visible text."""
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from lxml import etree as ET

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Add a title shape so 4.1.2 no-title doesn't fire (different rule)
    title_layout = prs.slide_layouts[5]  # title only
    # Use blank layout, just add a shape with title-like role via XML manipulation
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(1), Inches(2), Inches(1),
    )
    if descr:
        shape.element.nvSpPr.cNvPr.set('descr', descr)
    if has_text:
        shape.text_frame.text = "Click here"
    else:
        # Ensure no visible text
        shape.text_frame.text = ""

    # Inject an a:hlinkClick on cNvPr
    cNvPr = shape.element.nvSpPr.cNvPr
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
             'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
    hlink = ET.SubElement(cNvPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick',
                          nsmap=nsmap)
    hlink.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', '')
    hlink.set('action', 'ppaction://hlinksldjump')

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestPptx412ActionableShapeNames(unittest.TestCase):
    def test_fail_hyperlinked_shape_without_descr_or_text(self):
        data = _build_pptx_with_hyperlinked_shape(descr="", has_text=False)
        fs = PptxAnalyzer(data, "fail.pptx").analyze()
        findings = [
            f for f in _confirmed(fs, "4.1.2")
            if f.remediation_id and f.remediation_id.startswith("pptx_actionable_names_")
        ]
        self.assertEqual(len(findings), 1)
        self.assertIn("accessible name", findings[0].issue.lower())

    def test_pass_hyperlinked_shape_with_descr(self):
        data = _build_pptx_with_hyperlinked_shape(
            descr="Open project site", has_text=False
        )
        fs = PptxAnalyzer(data, "pass_descr.pptx").analyze()
        findings = [
            f for f in _confirmed(fs, "4.1.2")
            if f.remediation_id and f.remediation_id.startswith("pptx_actionable_names_")
        ]
        self.assertEqual(len(findings), 0)

    def test_pass_hyperlinked_shape_with_visible_text(self):
        data = _build_pptx_with_hyperlinked_shape(descr="", has_text=True)
        fs = PptxAnalyzer(data, "pass_text.pptx").analyze()
        findings = [
            f for f in _confirmed(fs, "4.1.2")
            if f.remediation_id and f.remediation_id.startswith("pptx_actionable_names_")
        ]
        self.assertEqual(len(findings), 0)


if __name__ == '__main__':
    unittest.main()
