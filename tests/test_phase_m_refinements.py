"""Phase M-refinements — per-format breadth additions.

R1: XLSX 2.4.2 — generic sheet names ("Sheet1", "Sheet 2", ...)
R2: XLSX 1.3.1 — large data sheets without frozen header rows
R3: DOCX 1.3.1 — tables without <w:tblHeader/> on the first row
R4: PPTX 4.1.2 — picture shapes with default cNvPr@name (POSSIBLE)
"""
import io
import unittest
import zipfile
from pathlib import Path

from openpyxl import Workbook

from wcag.analyzers.docx_analyzer import DocxAnalyzer
from wcag.analyzers.pptx_analyzer import PptxAnalyzer
from wcag.analyzers.xlsx_analyzer import XlsxAnalyzer


def _confirmed(fs, criterion: str, remediation_id_prefix: str):
    return [
        f for f in fs.confirmed_findings
        if f.criterion_id == criterion and (f.remediation_id or '').startswith(remediation_id_prefix)
    ]


def _possible(fs, criterion: str, remediation_id_prefix: str):
    return [
        f for f in fs.possible_findings
        if f.criterion_id == criterion and (f.remediation_id or '').startswith(remediation_id_prefix)
    ]


def _xlsx_bytes(wb: Workbook) -> bytes:
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# R1 — XLSX 2.4.2 generic sheet names
# ─────────────────────────────────────────────────────────────────────────────
class TestXlsxGenericSheetNames(unittest.TestCase):
    def test_fail_default_sheet_names(self):
        wb = Workbook()
        # Workbook starts with one default 'Sheet'; rename to match Excel default 'Sheet1'.
        wb.active.title = "Sheet1"
        wb.create_sheet("Sheet2")
        wb.create_sheet("Sheet 3")
        wb.properties.title = "Q4 Report"
        fs = XlsxAnalyzer(_xlsx_bytes(wb), "x.xlsx").analyze()
        findings = _confirmed(fs, "2.4.2", "xlsx_generic_sheet_names")
        self.assertEqual(len(findings), 1)
        self.assertIn("Sheet1", findings[0].evidence)

    def test_pass_descriptive_sheet_names(self):
        wb = Workbook()
        wb.active.title = "Revenue"
        wb.create_sheet("Headcount")
        wb.create_sheet("Forecast")
        wb.properties.title = "Q4 Report"
        fs = XlsxAnalyzer(_xlsx_bytes(wb), "x.xlsx").analyze()
        self.assertEqual(len(_confirmed(fs, "2.4.2", "xlsx_generic_sheet_names")), 0)


# ─────────────────────────────────────────────────────────────────────────────
# R2 — XLSX 1.3.1 unfrozen header on large sheet
# ─────────────────────────────────────────────────────────────────────────────
class TestXlsxUnfrozenHeader(unittest.TestCase):
    def _make_large_sheet(self, freeze: str | None) -> bytes:
        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        # Header row with text labels
        ws.append(["Region", "Revenue", "Headcount"])
        for i in range(25):
            ws.append([f"R{i}", 1000 + i, 10 + i])
        if freeze:
            ws.freeze_panes = freeze
        wb.properties.title = "Quarterly Data"
        return _xlsx_bytes(wb)

    def test_fail_no_freeze_panes(self):
        fs = XlsxAnalyzer(self._make_large_sheet(None), "x.xlsx").analyze()
        findings = _confirmed(fs, "1.3.1", "xlsx_unfrozen_header_")
        self.assertEqual(len(findings), 1)

    def test_pass_header_frozen(self):
        fs = XlsxAnalyzer(self._make_large_sheet("A2"), "x.xlsx").analyze()
        self.assertEqual(len(_confirmed(fs, "1.3.1", "xlsx_unfrozen_header_")), 0)

    def test_pass_small_sheet_skipped(self):
        wb = Workbook()
        ws = wb.active
        ws.title = "Tiny"
        ws.append(["Col1", "Col2"])
        for i in range(5):
            ws.append([f"R{i}", i])
        wb.properties.title = "Small"
        fs = XlsxAnalyzer(_xlsx_bytes(wb), "x.xlsx").analyze()
        self.assertEqual(len(_confirmed(fs, "1.3.1", "xlsx_unfrozen_header_")), 0)


# ─────────────────────────────────────────────────────────────────────────────
# R3 — DOCX 1.3.1 table header marking
# ─────────────────────────────────────────────────────────────────────────────
def _docx_with_table(table_xml: str) -> bytes:
    document_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    {table_xml}
    <w:sectPr><w:pgSz w:w="12240" w:h="15840"/><w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440" w:header="720" w:footer="720" w:gutter="0"/></w:sectPr>
  </w:body>
</w:document>"""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""")
        zf.writestr("_rels/.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""")
        zf.writestr("word/document.xml", document_xml)
    return buffer.getvalue()


class TestDocxTableHeaderMarking(unittest.TestCase):
    def test_fail_no_tblHeader(self):
        # 2-row table with no <w:tblHeader/>
        table = """<w:tbl>
          <w:tblPr><w:tblW w:type="auto" w:w="0"/></w:tblPr>
          <w:tr>
            <w:tc><w:p><w:r><w:t>Region</w:t></w:r></w:p></w:tc>
            <w:tc><w:p><w:r><w:t>Revenue</w:t></w:r></w:p></w:tc>
          </w:tr>
          <w:tr>
            <w:tc><w:p><w:r><w:t>North</w:t></w:r></w:p></w:tc>
            <w:tc><w:p><w:r><w:t>1000</w:t></w:r></w:p></w:tc>
          </w:tr>
        </w:tbl>"""
        fs = DocxAnalyzer(_docx_with_table(table), "t.docx").analyze()
        findings = _confirmed(fs, "1.3.1", "docx_table_header_marking")
        self.assertEqual(len(findings), 1)
        self.assertIn("Region", findings[0].evidence)

    def test_pass_with_tblHeader(self):
        table = """<w:tbl>
          <w:tblPr><w:tblW w:type="auto" w:w="0"/></w:tblPr>
          <w:tr>
            <w:trPr><w:tblHeader/></w:trPr>
            <w:tc><w:p><w:r><w:t>Region</w:t></w:r></w:p></w:tc>
            <w:tc><w:p><w:r><w:t>Revenue</w:t></w:r></w:p></w:tc>
          </w:tr>
          <w:tr>
            <w:tc><w:p><w:r><w:t>North</w:t></w:r></w:p></w:tc>
            <w:tc><w:p><w:r><w:t>1000</w:t></w:r></w:p></w:tc>
          </w:tr>
        </w:tbl>"""
        fs = DocxAnalyzer(_docx_with_table(table), "t.docx").analyze()
        self.assertEqual(len(_confirmed(fs, "1.3.1", "docx_table_header_marking")), 0)

    def test_pass_single_row_table_skipped(self):
        # 1-row table — exempt because it isn't really a data table
        table = """<w:tbl>
          <w:tblPr><w:tblW w:type="auto" w:w="0"/></w:tblPr>
          <w:tr>
            <w:tc><w:p><w:r><w:t>Just one row</w:t></w:r></w:p></w:tc>
          </w:tr>
        </w:tbl>"""
        fs = DocxAnalyzer(_docx_with_table(table), "t.docx").analyze()
        self.assertEqual(len(_confirmed(fs, "1.3.1", "docx_table_header_marking")), 0)


# ─────────────────────────────────────────────────────────────────────────────
# R4 — PPTX 4.1.2 generic picture names (POSSIBLE)
# ─────────────────────────────────────────────────────────────────────────────
def _build_pptx_with_picture_name(picture_name: str, descr: str = "") -> bytes:
    """Build a minimal PPTX with one slide containing one picture shape."""
    from pptx import Presentation
    from pptx.util import Inches
    import tempfile
    import os
    from PIL import Image

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    slide.shapes.title.text = "Slide title"
    # Need a real image file
    img_path = tempfile.mktemp(suffix=".png")
    Image.new("RGB", (50, 50), "red").save(img_path)
    try:
        pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1), Inches(1), Inches(1))
        pic.name = picture_name
        if descr:
            # Set alt text via the python-pptx descr API
            from pptx.oxml.ns import qn
            cNvPr = pic._element.nvPicPr.cNvPr
            cNvPr.set("descr", descr)
    finally:
        try:
            os.remove(img_path)
        except OSError:
            pass
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestPptxGenericPictureNames(unittest.TestCase):
    def test_fail_default_picture_name(self):
        data = _build_pptx_with_picture_name("Picture 1", descr="Logo")
        fs = PptxAnalyzer(data, "p.pptx").analyze()
        findings = _possible(fs, "4.1.2", "pptx_generic_picture_names_")
        self.assertEqual(len(findings), 1)
        self.assertIn("Picture 1", findings[0].evidence)

    def test_pass_descriptive_picture_name(self):
        data = _build_pptx_with_picture_name("CompanyLogo", descr="Logo")
        fs = PptxAnalyzer(data, "p.pptx").analyze()
        self.assertEqual(len(_possible(fs, "4.1.2", "pptx_generic_picture_names_")), 0)

    def test_pass_decorative_picture_skipped(self):
        # A decorative picture (descr unset is fine, but we mark decorative)
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
        import tempfile
        import os

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Slide"
        img_path = tempfile.mktemp(suffix=".png")
        Image.new("RGB", (50, 50), "red").save(img_path)
        try:
            pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1), Inches(1), Inches(1))
            pic.name = "Picture 1"
            # Mark decorative via OOXML extension
            from pptx.oxml.ns import qn
            from lxml import etree
            cNvPr = pic._element.nvPicPr.cNvPr
            extLst = etree.SubElement(cNvPr, qn("a:extLst"))
            ext = etree.SubElement(extLst, qn("a:ext"))
            ext.set("uri", "{C183D7F6-B498-43B3-948B-1728B52AA6E4}")
            adec = etree.SubElement(
                ext, "{http://schemas.microsoft.com/office/drawing/2017/decorative}decorative"
            )
            adec.set("val", "1")
        finally:
            try:
                os.remove(img_path)
            except OSError:
                pass
        buf = io.BytesIO()
        prs.save(buf)
        fs = PptxAnalyzer(buf.getvalue(), "p.pptx").analyze()
        self.assertEqual(len(_possible(fs, "4.1.2", "pptx_generic_picture_names_")), 0)


if __name__ == "__main__":
    unittest.main()
