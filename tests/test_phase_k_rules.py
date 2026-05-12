"""Phase K — cross-format coverage broadening.

K1: DOCX 1.4.10 — fixed-width tables wider than body
K2: PPTX 2.2.1 — auto-advance slide transitions
K3: XLSX 4.1.2 — form controls without objectName/altText
K4: PDF 1.3.5  — AcroForm fields without WCAG autocomplete tokens (POSSIBLE)
K5: HTML 2.5.7 — drag-only widgets (WCAG 2.2)

Each rule has at least one PASS and one FAIL test.
"""
import io
import unittest
import zipfile

from wcag.analyzers.docx_analyzer import DocxAnalyzer
from wcag.analyzers.pptx_analyzer import PptxAnalyzer
from wcag.analyzers.xlsx_analyzer import XlsxAnalyzer
from wcag.analyzers.pdf_analyzer import PdfAnalyzer
from wcag.analyzers.html_analyzer import HtmlAnalyzer


def _confirmed(fs, criterion: str):
    return [f for f in fs.confirmed_findings if f.criterion_id == criterion]


def _possible(fs, criterion: str):
    return [f for f in fs.possible_findings if f.criterion_id == criterion]


def _confirmed_with_id(fs, criterion: str, remediation_id: str):
    return [f for f in fs.confirmed_findings
            if f.criterion_id == criterion and f.remediation_id == remediation_id]


def _possible_with_id(fs, criterion: str, remediation_id: str):
    return [f for f in fs.possible_findings
            if f.criterion_id == criterion and f.remediation_id == remediation_id]


# ─────────────────────────────────────────────────────────────────────────────
# K1 — DOCX 1.4.10 fixed-width tables
# ─────────────────────────────────────────────────────────────────────────────
def _build_docx_with_table(tbl_width_dxa: int) -> bytes:
    """Build a minimal docx with a single table whose tblW is the given dxa."""
    from docx import Document
    doc = Document()
    doc.add_paragraph("Header paragraph.")
    table = doc.add_table(rows=2, cols=2)
    # python-docx writes tblW automatically. We patch the XML after save.
    buf = io.BytesIO()
    doc.save(buf)

    # Re-open zip, mutate document.xml, rewrite.
    src = zipfile.ZipFile(io.BytesIO(buf.getvalue()), 'r')
    parts = {n: src.read(n) for n in src.namelist()}
    src.close()

    doc_xml = parts['word/document.xml'].decode('utf-8')
    # Replace tblW value. python-docx's default may be type='auto' w='0'.
    # We force type='dxa' w='<our value>'.
    import re as _re
    new_tblW = f'<w:tblW w:w="{tbl_width_dxa}" w:type="dxa"/>'
    if '<w:tblW' in doc_xml:
        doc_xml = _re.sub(r'<w:tblW[^/]*/>', new_tblW, doc_xml, count=1)
    else:
        # Insert into <w:tblPr>
        doc_xml = doc_xml.replace('<w:tblPr>', f'<w:tblPr>{new_tblW}', 1)
    parts['word/document.xml'] = doc_xml.encode('utf-8')

    out = io.BytesIO()
    with zipfile.ZipFile(out, 'w', zipfile.ZIP_DEFLATED) as zf:
        for n, data in parts.items():
            zf.writestr(n, data)
    return out.getvalue()


class TestDocx1410FixedWidthTables(unittest.TestCase):
    def test_pass_table_within_body(self):
        # 6 inches wide in dxa = 8640. Default body ~9360. Should NOT fire.
        data = _build_docx_with_table(tbl_width_dxa=8640)
        fs = DocxAnalyzer(data, 'fits.docx').analyze()
        self.assertEqual(
            len(_confirmed_with_id(fs, "1.4.10", "docx_fixed_width_tables")), 0
        )

    def test_fail_table_wider_than_body(self):
        # 11 inches wide = 15840 dxa, way wider than ~9360 dxa body.
        data = _build_docx_with_table(tbl_width_dxa=15840)
        fs = DocxAnalyzer(data, 'wide.docx').analyze()
        findings = _confirmed_with_id(fs, "1.4.10", "docx_fixed_width_tables")
        self.assertEqual(len(findings), 1)
        self.assertIn("table", findings[0].issue.lower())


# ─────────────────────────────────────────────────────────────────────────────
# K2 — PPTX 2.2.1 auto-advance slides
# ─────────────────────────────────────────────────────────────────────────────
def _build_pptx_with_advance(advance_ms: int = None) -> bytes:
    """Build a minimal pptx; if advance_ms is set, inject <p:transition advTm='N'/>
    onto slide1.xml."""
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Sample Slide"
    buf = io.BytesIO()
    prs.save(buf)

    if advance_ms is None:
        return buf.getvalue()

    src = zipfile.ZipFile(io.BytesIO(buf.getvalue()), 'r')
    parts = {n: src.read(n) for n in src.namelist()}
    src.close()

    slide_xml = parts['ppt/slides/slide1.xml'].decode('utf-8')
    transition = (
        f'<p:transition xmlns:p="http://schemas.openxmlformats.org/'
        f'presentationml/2006/main" advTm="{advance_ms}"/>'
    )
    # Insert before </p:sld>
    if '</p:sld>' in slide_xml:
        slide_xml = slide_xml.replace('</p:sld>', f'{transition}</p:sld>', 1)
    parts['ppt/slides/slide1.xml'] = slide_xml.encode('utf-8')

    out = io.BytesIO()
    with zipfile.ZipFile(out, 'w', zipfile.ZIP_DEFLATED) as zf:
        for n, data in parts.items():
            zf.writestr(n, data)
    return out.getvalue()


class TestPptx221AutoAdvance(unittest.TestCase):
    def test_pass_no_auto_advance(self):
        data = _build_pptx_with_advance(advance_ms=None)
        fs = PptxAnalyzer(data, 'manual.pptx').analyze()
        self.assertEqual(
            len(_confirmed_with_id(fs, "2.2.1", "pptx_auto_advance_slides")), 0
        )

    def test_fail_auto_advance_5_seconds(self):
        data = _build_pptx_with_advance(advance_ms=5000)
        fs = PptxAnalyzer(data, 'auto.pptx').analyze()
        findings = _confirmed_with_id(fs, "2.2.1", "pptx_auto_advance_slides")
        self.assertEqual(len(findings), 1)
        self.assertIn("auto-advance", findings[0].issue.lower())


# ─────────────────────────────────────────────────────────────────────────────
# K3 — XLSX 4.1.2 form controls without name
# ─────────────────────────────────────────────────────────────────────────────
def _build_xlsx_with_ctrl_props(object_name: str = None, alt_text: str = None) -> bytes:
    """Build a minimal xlsx with one xl/ctrlProps/ctrlProp1.xml file."""
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws['A1'] = "Sample"
    buf = io.BytesIO()
    wb.save(buf)

    src = zipfile.ZipFile(io.BytesIO(buf.getvalue()), 'r')
    parts = {n: src.read(n) for n in src.namelist()}
    src.close()

    attrs = []
    if object_name is not None:
        attrs.append(f'objectName="{object_name}"')
    if alt_text is not None:
        attrs.append(f'altText="{alt_text}"')
    attrs.append('objectType="CheckBox"')
    ctrl_xml = (
        f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        f'<formControlPr xmlns="http://schemas.microsoft.com/office/spreadsheetml/2009/9/main" '
        f'{" ".join(attrs)}/>'
    )
    parts['xl/ctrlProps/ctrlProp1.xml'] = ctrl_xml.encode('utf-8')

    out = io.BytesIO()
    with zipfile.ZipFile(out, 'w', zipfile.ZIP_DEFLATED) as zf:
        for n, data in parts.items():
            zf.writestr(n, data)
    return out.getvalue()


class TestXlsx412FormControls(unittest.TestCase):
    def test_pass_control_with_object_name(self):
        data = _build_xlsx_with_ctrl_props(object_name="OptIn-Newsletter")
        fs = XlsxAnalyzer(data, 'named.xlsx').analyze()
        self.assertEqual(
            len(_confirmed_with_id(fs, "4.1.2", "xlsx_form_control_names")), 0
        )

    def test_pass_control_with_alt_text_only(self):
        data = _build_xlsx_with_ctrl_props(alt_text="Subscribe to newsletter")
        fs = XlsxAnalyzer(data, 'alt_only.xlsx').analyze()
        self.assertEqual(
            len(_confirmed_with_id(fs, "4.1.2", "xlsx_form_control_names")), 0
        )

    def test_fail_control_with_no_name(self):
        # Neither objectName nor altText — must fire.
        data = _build_xlsx_with_ctrl_props()
        fs = XlsxAnalyzer(data, 'unnamed.xlsx').analyze()
        findings = _confirmed_with_id(fs, "4.1.2", "xlsx_form_control_names")
        self.assertEqual(len(findings), 1)
        self.assertIn("form control", findings[0].issue.lower())


# ─────────────────────────────────────────────────────────────────────────────
# K4 — PDF 1.3.5 input purpose (POSSIBLE)
# ─────────────────────────────────────────────────────────────────────────────
def _build_pdf_with_form_field(field_name: str, tooltip: str = None) -> bytes:
    """Build a minimal PDF with a single AcroForm text field."""
    import pikepdf
    from pikepdf import Pdf, Dictionary, Array, Name, String

    pdf = Pdf.new()
    pdf.add_blank_page(page_size=(612, 792))

    field = pdf.make_indirect(Dictionary({
        '/Type': Name('/Annot'),
        '/Subtype': Name('/Widget'),
        '/FT': Name('/Tx'),
        '/T': String(field_name),
        '/Rect': Array([100, 100, 300, 130]),
    }))
    if tooltip:
        field['/TU'] = String(tooltip)

    acroform = Dictionary({'/Fields': Array([field])})
    pdf.Root['/AcroForm'] = acroform

    buf = io.BytesIO()
    pdf.save(buf)
    return buf.getvalue()


class TestPdf135InputPurpose(unittest.TestCase):
    def test_pass_field_uses_canonical_token(self):
        data = _build_pdf_with_form_field(field_name="email")
        fs = PdfAnalyzer(data, 'good.pdf').analyze()
        self.assertEqual(
            len(_possible_with_id(fs, "1.3.5", "pdf_input_purpose")), 0
        )

    def test_pass_unknown_field_not_flagged(self):
        # A field name with no recognised purpose hint should NOT fire.
        data = _build_pdf_with_form_field(field_name="comments")
        fs = PdfAnalyzer(data, 'comments.pdf').analyze()
        self.assertEqual(
            len(_possible_with_id(fs, "1.3.5", "pdf_input_purpose")), 0
        )

    def test_fail_field_named_phone_not_canonical(self):
        # 'phone' is a hint for 'tel' but is not a canonical token.
        data = _build_pdf_with_form_field(field_name="phone_number")
        fs = PdfAnalyzer(data, 'phone.pdf').analyze()
        findings = _possible_with_id(fs, "1.3.5", "pdf_input_purpose")
        self.assertEqual(len(findings), 1)
        self.assertIn("autocomplete", findings[0].evidence.lower())


# ─────────────────────────────────────────────────────────────────────────────
# K5 — HTML 2.5.7 dragging movements (WCAG 2.2)
# ─────────────────────────────────────────────────────────────────────────────
class TestHtml257DraggingMovements(unittest.TestCase):
    def test_pass_no_drag_widgets(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><button>Move up</button><button>Move down</button></body></html>"""
        fs = HtmlAnalyzer(html.encode("utf-8"), 'pass.html').analyze()
        self.assertEqual(
            len(_confirmed_with_id(fs, "2.5.7", "html_dragging_movements")), 0
        )
        self.assertEqual(
            len(_possible_with_id(fs, "2.5.7", "html_dragging_movements")), 0
        )

    def test_fail_draggable_with_no_alternative(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body>
          <ul>
            <li draggable='true'>Item 1</li>
            <li draggable='true'>Item 2</li>
          </ul>
        </body></html>"""
        fs = HtmlAnalyzer(html.encode("utf-8"), 'drag.html').analyze()
        findings = _confirmed_with_id(fs, "2.5.7", "html_dragging_movements")
        self.assertEqual(len(findings), 1)
        self.assertIn("draggable", findings[0].issue.lower())

    def test_possible_when_page_has_move_buttons(self):
        # Draggable widget present, but page also has 'Move up' button —
        # downgrade to POSSIBLE.
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body>
          <ul>
            <li draggable='true'>Item 1</li>
          </ul>
          <button>Move up</button>
        </body></html>"""
        fs = HtmlAnalyzer(html.encode("utf-8"), 'mixed.html').analyze()
        confirmed = _confirmed_with_id(fs, "2.5.7", "html_dragging_movements")
        possible = _possible_with_id(fs, "2.5.7", "html_dragging_movements")
        self.assertEqual(len(confirmed), 0)
        self.assertEqual(len(possible), 1)


if __name__ == '__main__':
    unittest.main()
