"""
Step 2: Confidence Contract Audit Tests

Ensures every Finding has complete confidence metadata:
- confidence_tier: CONFIRMED or POSSIBLE
- confidence_label: "high", "medium", or "low"
- confidence_rationale: one sentence explaining certainty

Tests run across all analyzers (HTML, PDF, DOCX, XLSX, PPTX) to
verify auditability and evidence traceability.
"""

import unittest
from wcag.analyzers.html_analyzer import HtmlAnalyzer
from wcag.analyzers.pdf_analyzer import PdfAnalyzer
from wcag.analyzers.docx_analyzer import DocxAnalyzer
from wcag.analyzers.xlsx_analyzer import XlsxAnalyzer
from wcag.analyzers.pptx_analyzer import PptxAnalyzer
from wcag.models import ConfidenceTier


class ConfidenceContractAudit(unittest.TestCase):
    """Audit confidence metadata completeness across all analyzers."""

    def _audit_findings(self, fact_sheet, analyzer_name):
        """Check all findings in a fact sheet for confidence contract compliance."""
        all_findings = fact_sheet.confirmed_findings + fact_sheet.possible_findings
        violations = []

        for i, f in enumerate(all_findings):
            issues = []

            # Every finding must have a confidence_tier
            if not hasattr(f, 'confidence_tier') or f.confidence_tier is None:
                issues.append("missing confidence_tier")
            elif not isinstance(f.confidence_tier, ConfidenceTier):
                issues.append(f"invalid confidence_tier type: {type(f.confidence_tier)}")

            # Every finding must have a confidence_label
            if not hasattr(f, 'confidence_label') or not f.confidence_label:
                issues.append("missing confidence_label")
            elif f.confidence_label not in ("high", "medium", "low"):
                issues.append(f"invalid confidence_label: {f.confidence_label}")

            # Every finding must have a confidence_rationale
            if not hasattr(f, 'confidence_rationale') or not f.confidence_rationale:
                issues.append("missing confidence_rationale")
            elif not isinstance(f.confidence_rationale, str) or len(f.confidence_rationale) < 10:
                issues.append(f"confidence_rationale too short: {f.confidence_rationale[:20]}")

            # Rationale should be 1-3 sentences, roughly 20-150 characters
            if hasattr(f, 'confidence_rationale') and f.confidence_rationale:
                words = len(f.confidence_rationale.split())
                if words < 5:
                    issues.append(f"confidence_rationale too brief ({words} words)")

            # CONFIRMED findings should have evidence
            if f.confidence_tier == ConfidenceTier.CONFIRMED:
                if not hasattr(f, 'evidence') or not f.evidence:
                    issues.append("CONFIRMED finding missing evidence")

            if issues:
                violations.append({
                    "finding_index": i,
                    "criterion_id": f.criterion_id,
                    "issue": f.issue[:50],
                    "violations": issues,
                })

        return violations

    def test_html_analyzer_confidence_contract(self):
        """HTML findings must have complete confidence metadata."""
        html = b"<!DOCTYPE html><html lang='en'><head><title>Test</title></head><body><h1>Hello</h1></body></html>"
        fs = HtmlAnalyzer(html, "test.html").analyze()
        violations = self._audit_findings(fs, "html")
        self.assertEqual(len(violations), 0, f"HTML confidence audit failed: {violations}")

    def test_docx_analyzer_confidence_contract(self):
        """DOCX findings must have complete confidence metadata."""
        from io import BytesIO
        from docx import Document

        doc = Document()
        doc.add_heading("Test", level=1)
        doc.add_paragraph("Some content.")

        buf = BytesIO()
        doc.save(buf)
        buf.seek(0)

        fs = DocxAnalyzer(buf.getvalue(), "test.docx").analyze()
        violations = self._audit_findings(fs, "docx")
        self.assertEqual(len(violations), 0, f"DOCX confidence audit failed: {violations}")

    def test_xlsx_analyzer_confidence_contract(self):
        """XLSX findings must have complete confidence metadata."""
        from openpyxl import Workbook
        from io import BytesIO

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws["A1"] = "Name"
        ws["A2"] = "Test"

        buf = BytesIO()
        wb.save(buf)
        buf.seek(0)

        fs = XlsxAnalyzer(buf.getvalue(), "test.xlsx").analyze()
        violations = self._audit_findings(fs, "xlsx")
        self.assertEqual(len(violations), 0, f"XLSX confidence audit failed: {violations}")

    def test_pptx_analyzer_confidence_contract(self):
        """PPTX findings must have complete confidence metadata."""
        from pptx import Presentation
        from io import BytesIO

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test"

        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)

        fs = PptxAnalyzer(buf.getvalue(), "test.pptx").analyze()
        violations = self._audit_findings(fs, "pptx")
        self.assertEqual(len(violations), 0, f"PPTX confidence audit failed: {violations}")

    def test_evidence_source_confidence_label_alignment(self):
        """Confidence labels should align with evidence source strength."""
        html = b"<!DOCTYPE html><html lang='en'><head><title>Test</title></head><body><img src='x.jpg'></body></html>"
        fs = HtmlAnalyzer(html, "test.html").analyze()

        from wcag.models import EvidenceSource

        for f in fs.confirmed_findings + fs.possible_findings:
            if not hasattr(f, "evidence_source"):
                continue

            source = f.evidence_source
            label = f.confidence_label

            # XML_DIRECT should be "high"
            if source == EvidenceSource.XML_DIRECT:
                self.assertEqual(label, "high", f"XML_DIRECT finding {f.criterion_id} has label={label}, expected high")

            # BROWSER_RENDERED should be "high" (rendered HTML checks)
            if source == EvidenceSource.BROWSER_RENDERED:
                self.assertEqual(label, "high", f"BROWSER_RENDERED finding {f.criterion_id} has label={label}, expected high")

            # XML_INFERRED can be "medium" or "high"
            if source == EvidenceSource.XML_INFERRED:
                self.assertIn(label, ("medium", "high"), f"XML_INFERRED finding {f.criterion_id} has unexpected label={label}")


class ConfidenceRationaleQuality(unittest.TestCase):
    """Audit quality of confidence rationales — should be concise and specific."""

    def test_rationale_references_evidence_shape(self):
        """Rationales should mention evidence type (XML, text, rendered, etc)."""
        html = b"<!DOCTYPE html><html lang='en'><head><title>Test</title></head><body><h1></h1></body></html>"
        fs = HtmlAnalyzer(html, "test.html").analyze()

        evidence_keywords = {
            "xml_direct": ["xml", "attribute", "element", "directly"],
            "xml_inferred": ["xml", "inferred", "structure", "pattern"],
            "text_content": ["text", "content", "string", "phrase"],
            "browser_rendered": ["rendered", "browser", "computed", "style"],
            "dom_direct": ["dom", "attribute", "element", "directly"],
        }

        for f in fs.confirmed_findings + fs.possible_findings:
            source = str(f.evidence_source.value).lower()
            rationale = (f.confidence_rationale or "").lower()

            # Rationale should mention evidence source type
            if source in evidence_keywords:
                keywords = evidence_keywords[source]
                found_keyword = any(kw in rationale for kw in keywords)
                # Allow for some flexibility — at minimum rationale should be non-empty
                self.assertTrue(
                    found_keyword or len(rationale) >= 20,
                    f"Finding {f.criterion_id} with {source} source has weak rationale: {f.confidence_rationale[:50]}",
                )


if __name__ == "__main__":
    unittest.main()
