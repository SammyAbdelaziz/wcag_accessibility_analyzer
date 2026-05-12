"""Phase B tests — WCAG 1.4.11 Non-text Contrast across HTML/DOCX/PPTX/XLSX.

Each analyzer is exercised with one PASS fixture (>= 3:1) and one FAIL
fixture (< 3:1). The shared evaluator is unit-tested separately in
test_non_text_contrast.py.
"""
import io
import unittest
import zipfile

from wcag.analyzers.html_analyzer import HtmlAnalyzer
from wcag.analyzers.docx_analyzer import DocxAnalyzer
from wcag.analyzers.pptx_analyzer import PptxAnalyzer
from wcag.analyzers.xlsx_analyzer import XlsxAnalyzer


def _findings_1_4_11(fs):
    return [f for f in fs.confirmed_findings if f.criterion_id == "1.4.11"]


# ─────────────────────────────────────────────────────────────────────────────
# HTML
# ─────────────────────────────────────────────────────────────────────────────
class TestHtml1411(unittest.TestCase):
    def _analyze(self, html: str):
        return HtmlAnalyzer(html.encode('utf-8'), 'ntc.html').analyze()

    def test_pass_high_contrast_button_border(self):
        # Black border on white background — 21:1 — should NOT trigger 1.4.11
        html = """<!DOCTYPE html><html lang='en'><head><title>t</title></head><body>
        <button style='border: 2px solid #000000; background:#FFFFFF;'>Submit</button>
        </body></html>"""
        fs = self._analyze(html)
        self.assertEqual(len(_findings_1_4_11(fs)), 0)

    def test_fail_low_contrast_button_border(self):
        # Light gray (#CCCCCC) border on white — 1.61:1 — should trigger 1.4.11
        html = """<!DOCTYPE html><html lang='en'><head><title>t</title></head><body>
        <button style='border: 2px solid #CCCCCC; background:#FFFFFF;'>Submit</button>
        </body></html>"""
        fs = self._analyze(html)
        findings = _findings_1_4_11(fs)
        self.assertGreater(len(findings), 0)


# ─────────────────────────────────────────────────────────────────────────────
# PPTX
# ─────────────────────────────────────────────────────────────────────────────
def _build_pptx_with_shape(line_hex: str, fill_hex: str) -> bytes:
    """Build a minimal pptx with one rectangle shape having explicit
    line color and explicit fill color."""
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    blank = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank)
    # Add a title placeholder so 4.1.2 doesn't fire excessively
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(0.5))
    txBox.text_frame.text = "Slide One"
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1), Inches(1), Inches(2), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    shape.line.color.rgb = RGBColor.from_string(line_hex)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestPptx1411(unittest.TestCase):
    def test_pass_high_contrast_shape(self):
        # Black outline on white fill — 21:1
        data = _build_pptx_with_shape('000000', 'FFFFFF')
        fs = PptxAnalyzer(data, 'ntc.pptx').analyze()
        self.assertEqual(len(_findings_1_4_11(fs)), 0)

    def test_fail_low_contrast_shape(self):
        # Light gray outline on white fill — 1.61:1
        data = _build_pptx_with_shape('CCCCCC', 'FFFFFF')
        fs = PptxAnalyzer(data, 'ntc.pptx').analyze()
        self.assertGreater(len(_findings_1_4_11(fs)), 0)


# ─────────────────────────────────────────────────────────────────────────────
# XLSX
# ─────────────────────────────────────────────────────────────────────────────
def _build_xlsx_with_styled_cell(border_argb: str, fill_argb: str) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import PatternFill, Border, Side, Color

    wb = Workbook()
    ws = wb.active
    ws['A1'] = "Header"
    side = Side(style='thin', color=Color(rgb=border_argb))
    ws['A1'].border = Border(left=side, right=side, top=side, bottom=side)
    ws['A1'].fill = PatternFill(patternType='solid', fgColor=Color(rgb=fill_argb))
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


class TestXlsx1411(unittest.TestCase):
    def test_pass_high_contrast_cell(self):
        # Border #000000 on fill #FFFFFF — 21:1
        data = _build_xlsx_with_styled_cell('FF000000', 'FFFFFFFF')
        fs = XlsxAnalyzer(data, 'ntc.xlsx').analyze()
        self.assertEqual(len(_findings_1_4_11(fs)), 0)

    def test_fail_low_contrast_cell(self):
        # Border #CCCCCC on fill #FFFFFF — 1.61:1
        data = _build_xlsx_with_styled_cell('FFCCCCCC', 'FFFFFFFF')
        fs = XlsxAnalyzer(data, 'ntc.xlsx').analyze()
        self.assertGreater(len(_findings_1_4_11(fs)), 0)


# ─────────────────────────────────────────────────────────────────────────────
# DOCX
# ─────────────────────────────────────────────────────────────────────────────
def _build_docx_with_styled_table(border_hex: str, fill_hex: str) -> bytes:
    """Build a docx with a 2x2 table whose first cell has explicit
    border color + shading fill, by post-processing the underlying XML."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("test")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A"
    buf = io.BytesIO()
    doc.save(buf)
    raw = buf.getvalue()

    # Inject explicit tcBorders + shd into the first <w:tc>.
    out = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(raw), 'r') as zin:
        with zipfile.ZipFile(out, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == 'word/document.xml':
                    text = data.decode('utf-8')
                    # Inject a <w:tcPr> with explicit borders + shd into the first cell
                    inject = (
                        '<w:tcPr xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
                        '<w:tcBorders>'
                        f'<w:top w:val="single" w:sz="4" w:space="0" w:color="{border_hex}"/>'
                        f'<w:left w:val="single" w:sz="4" w:space="0" w:color="{border_hex}"/>'
                        f'<w:bottom w:val="single" w:sz="4" w:space="0" w:color="{border_hex}"/>'
                        f'<w:right w:val="single" w:sz="4" w:space="0" w:color="{border_hex}"/>'
                        '</w:tcBorders>'
                        f'<w:shd w:val="clear" w:color="auto" w:fill="{fill_hex}"/>'
                        '</w:tcPr>'
                    )
                    # First <w:tc> opens; insert tcPr right after it
                    text = text.replace('<w:tc>', '<w:tc>' + inject, 1)
                    data = text.encode('utf-8')
                zout.writestr(item, data)
    return out.getvalue()


class TestDocx1411(unittest.TestCase):
    def test_pass_high_contrast_cell(self):
        # Border #000000 on fill #FFFFFF — 21:1
        data = _build_docx_with_styled_table('000000', 'FFFFFF')
        fs = DocxAnalyzer(data, 'ntc.docx').analyze()
        self.assertEqual(len(_findings_1_4_11(fs)), 0)

    def test_fail_low_contrast_cell(self):
        # Border #CCCCCC on fill #FFFFFF — 1.61:1
        data = _build_docx_with_styled_table('CCCCCC', 'FFFFFF')
        fs = DocxAnalyzer(data, 'ntc.docx').analyze()
        self.assertGreater(len(_findings_1_4_11(fs)), 0)


if __name__ == '__main__':
    unittest.main()
