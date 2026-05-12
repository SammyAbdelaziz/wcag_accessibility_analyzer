"""
Phase A — New Rule Coverage Unit Tests

Deterministic pass + fail tests for each Phase A rule added across 5 analyzers.
"""

import unittest
import io
import zipfile
from wcag.analyzers.html_analyzer import HtmlAnalyzer
from wcag.analyzers.docx_analyzer import DocxAnalyzer
from wcag.analyzers.pptx_analyzer import PptxAnalyzer
from wcag.analyzers.xlsx_analyzer import XlsxAnalyzer


def _make_html(body: str, head: str = "<title>Test</title>") -> bytes:
    return (
        f"<!DOCTYPE html><html lang='en'><head>{head}</head><body>{body}</body></html>"
    ).encode("utf-8")


def _findings_for(fs, criterion: str, remediation_id: str = None):
    all_findings = fs.confirmed_findings + fs.possible_findings
    matches = [f for f in all_findings if f.criterion_id == criterion]
    if remediation_id:
        matches = [f for f in matches if f.remediation_id == remediation_id]
    return matches


# =============================================================================
# HTML 2.4.1 — Bypass Blocks
# =============================================================================

class TestHtml241BypassBlocks(unittest.TestCase):
    def test_nav_with_main_landmark_passes(self):
        body = "<nav><a href='/a'>A</a></nav><main>Content</main>"
        fs = HtmlAnalyzer(_make_html(body), "ok.html").analyze()
        self.assertEqual(len(_findings_for(fs, "2.4.1", "html_bypass_blocks")), 0)

    def test_nav_with_skip_link_passes(self):
        body = "<a href='#main' class='skip'>Skip to main</a><nav><a href='/a'>A</a></nav><div id='main'>X</div>"
        fs = HtmlAnalyzer(_make_html(body), "ok.html").analyze()
        self.assertEqual(len(_findings_for(fs, "2.4.1", "html_bypass_blocks")), 0)

    def test_nav_without_bypass_fails(self):
        body = "<nav><a href='/a'>A</a><a href='/b'>B</a></nav><div>Content</div>"
        fs = HtmlAnalyzer(_make_html(body), "bad.html").analyze()
        self.assertGreater(len(_findings_for(fs, "2.4.1", "html_bypass_blocks")), 0)

    def test_no_nav_no_finding(self):
        body = "<p>Just content</p>"
        fs = HtmlAnalyzer(_make_html(body), "ok.html").analyze()
        self.assertEqual(len(_findings_for(fs, "2.4.1", "html_bypass_blocks")), 0)


# =============================================================================
# HTML 1.3.5 — Input Purpose (autocomplete)
# =============================================================================

class TestHtml135InputPurpose(unittest.TestCase):
    def test_email_input_with_autocomplete_passes(self):
        body = "<form><input type='email' name='email' autocomplete='email'></form>"
        fs = HtmlAnalyzer(_make_html(body), "ok.html").analyze()
        self.assertEqual(len(_findings_for(fs, "1.3.5", "html_input_purpose")), 0)

    def test_email_without_autocomplete_fails(self):
        body = "<form><input type='email' name='email'></form>"
        fs = HtmlAnalyzer(_make_html(body), "bad.html").analyze()
        self.assertGreater(len(_findings_for(fs, "1.3.5", "html_input_purpose")), 0)


# =============================================================================
# HTML 3.3.2 — Labels or Instructions
# =============================================================================

class TestHtml332Labels(unittest.TestCase):
    def test_input_with_label_passes(self):
        body = "<label for='e'>Email</label><input type='email' id='e'>"
        fs = HtmlAnalyzer(_make_html(body), "ok.html").analyze()
        self.assertEqual(len(_findings_for(fs, "3.3.2", "html_input_labels")), 0)

    def test_input_with_aria_label_passes(self):
        body = "<input type='text' aria-label='Search'>"
        fs = HtmlAnalyzer(_make_html(body), "ok.html").analyze()
        self.assertEqual(len(_findings_for(fs, "3.3.2", "html_input_labels")), 0)

    def test_unlabeled_input_fails(self):
        body = "<input type='text' name='nameless'>"
        fs = HtmlAnalyzer(_make_html(body), "bad.html").analyze()
        self.assertGreater(len(_findings_for(fs, "3.3.2", "html_input_labels")), 0)


# =============================================================================
# HTML 2.4.5 — Multiple Ways
# =============================================================================

class TestHtml245MultipleWays(unittest.TestCase):
    def test_with_nav_passes(self):
        body = "<nav><a href='/a'>A</a></nav><p>X</p>"
        fs = HtmlAnalyzer(_make_html(body), "ok.html").analyze()
        self.assertEqual(len(_findings_for(fs, "2.4.5", "html_multiple_ways")), 0)

    def test_with_search_passes(self):
        body = "<input type='search'><p>X</p>"
        fs = HtmlAnalyzer(_make_html(body), "ok.html").analyze()
        self.assertEqual(len(_findings_for(fs, "2.4.5", "html_multiple_ways")), 0)

    def test_no_nav_no_search_no_sitemap_flagged(self):
        # Page links to multiple internal pages → looks like a multi-page site
        # without nav/search/sitemap → should trigger 2.4.5.
        body = """
        <p>Plain content</p>
        <a href='/about'>About</a>
        <a href='/products'>Products</a>
        <a href='/contact'>Contact</a>
        <a href='/blog'>Blog</a>
        """
        fs = HtmlAnalyzer(_make_html(body), "bad.html").analyze()
        self.assertGreater(len(_findings_for(fs, "2.4.5", "html_multiple_ways")), 0)


# =============================================================================
# HTML 1.4.4 — Resize Text (viewport)
# =============================================================================

class TestHtml144ViewportResize(unittest.TestCase):
    def test_no_viewport_no_finding(self):
        fs = HtmlAnalyzer(_make_html("<p>X</p>"), "ok.html").analyze()
        self.assertEqual(len(_findings_for(fs, "1.4.4", "html_viewport_zoom")), 0)

    def test_proper_viewport_passes(self):
        head = "<title>T</title><meta name='viewport' content='width=device-width, initial-scale=1'>"
        fs = HtmlAnalyzer(_make_html("<p>X</p>", head=head), "ok.html").analyze()
        self.assertEqual(len(_findings_for(fs, "1.4.4", "html_viewport_zoom")), 0)

    def test_user_scalable_no_fails(self):
        head = "<title>T</title><meta name='viewport' content='width=device-width, user-scalable=no'>"
        fs = HtmlAnalyzer(_make_html("<p>X</p>", head=head), "bad.html").analyze()
        self.assertGreater(len(_findings_for(fs, "1.4.4", "html_viewport_zoom")), 0)

    def test_max_scale_too_low_fails(self):
        head = "<title>T</title><meta name='viewport' content='width=device-width, maximum-scale=1.0'>"
        fs = HtmlAnalyzer(_make_html("<p>X</p>", head=head), "bad.html").analyze()
        self.assertGreater(len(_findings_for(fs, "1.4.4", "html_viewport_zoom")), 0)


# =============================================================================
# DOCX 2.4.5 — Multiple Ways (TOC/bookmarks)
# =============================================================================

class TestDocx245MultipleWays(unittest.TestCase):
    def test_short_doc_no_finding(self):
        from docx import Document
        doc = Document()
        doc.add_heading("H1", level=1)
        doc.add_paragraph("body")
        buf = io.BytesIO()
        doc.save(buf)
        fs = DocxAnalyzer(buf.getvalue(), "short.docx").analyze()
        self.assertEqual(len(_findings_for(fs, "2.4.5", "docx_multiple_ways")), 0)

    def test_long_doc_without_toc_flagged(self):
        from docx import Document
        doc = Document()
        for i in range(6):
            doc.add_heading(f"Section {i+1}", level=1)
            doc.add_paragraph(f"Body of section {i+1}.")
        buf = io.BytesIO()
        doc.save(buf)
        fs = DocxAnalyzer(buf.getvalue(), "long.docx").analyze()
        self.assertGreater(len(_findings_for(fs, "2.4.5", "docx_multiple_ways")), 0)


# =============================================================================
# PPTX 2.4.6 — Slide Title Quality
# =============================================================================

class TestPptx246TitleQuality(unittest.TestCase):
    def test_unique_titles_pass(self):
        from pptx import Presentation
        prs = Presentation()
        prs.core_properties.title = "Deck"
        for i, t in enumerate(["Overview", "Q1 Sales", "Q2 Sales"]):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = t
        buf = io.BytesIO()
        prs.save(buf)
        fs = PptxAnalyzer(buf.getvalue(), "good.pptx").analyze()
        self.assertEqual(len(_findings_for(fs, "2.4.6", "pptx_title_duplicates")), 0)

    def test_duplicate_titles_flagged(self):
        from pptx import Presentation
        prs = Presentation()
        prs.core_properties.title = "Deck"
        for t in ["Overview", "Sales", "Sales"]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = t
        buf = io.BytesIO()
        prs.save(buf)
        fs = PptxAnalyzer(buf.getvalue(), "dup.pptx").analyze()
        self.assertGreater(len(_findings_for(fs, "2.4.6", "pptx_title_duplicates")), 0)


# =============================================================================
# XLSX 2.4.6 — Header Row Quality
# =============================================================================

class TestXlsx246HeaderQuality(unittest.TestCase):
    def _build(self, headers):
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        for col_idx, h in enumerate(headers, 1):
            ws.cell(row=1, column=col_idx, value=h)
            ws.cell(row=2, column=col_idx, value=f"data{col_idx}")
            ws.cell(row=3, column=col_idx, value=f"data{col_idx}b")
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()

    def test_descriptive_headers_pass(self):
        data = self._build(["Name", "Age", "City"])
        fs = XlsxAnalyzer(data, "ok.xlsx").analyze()
        # No header quality finding for descriptive headers
        matches = [f for f in fs.possible_findings
                   if f.criterion_id == "2.4.6" and "header_quality" in (f.remediation_id or "")]
        self.assertEqual(len(matches), 0)

    def test_generic_headers_flagged(self):
        data = self._build(["Column1", "Column2", "Column3"])
        fs = XlsxAnalyzer(data, "bad.xlsx").analyze()
        matches = [f for f in fs.possible_findings
                   if f.criterion_id == "2.4.6" and "header_quality" in (f.remediation_id or "")]
        self.assertGreater(len(matches), 0)

    def test_duplicate_headers_flagged(self):
        data = self._build(["Sales", "Sales", "Sales"])
        fs = XlsxAnalyzer(data, "bad.xlsx").analyze()
        matches = [f for f in fs.possible_findings
                   if f.criterion_id == "2.4.6" and "header_quality" in (f.remediation_id or "")]
        self.assertGreater(len(matches), 0)


# =============================================================================
# XLSX 3.1.2 — Language of Parts
# =============================================================================

class TestXlsx312LanguageOfParts(unittest.TestCase):
    def test_no_lang_attrs_passes(self):
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws['A1'] = "Hello"
        buf = io.BytesIO()
        wb.save(buf)
        fs = XlsxAnalyzer(buf.getvalue(), "ok.xlsx").analyze()
        matches = [f for f in fs.possible_findings
                   if f.criterion_id == "3.1.2" and f.remediation_id == "xlsx_language_of_parts"]
        self.assertEqual(len(matches), 0)

    def test_foreign_lang_attrs_flagged(self):
        # Build a minimal xlsx then inject xml:lang into the inline string
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws['A1'] = "Hello"
        buf = io.BytesIO()
        wb.save(buf)
        original = buf.getvalue()

        out_buf = io.BytesIO()
        with zipfile.ZipFile(io.BytesIO(original), 'r') as zin:
            with zipfile.ZipFile(out_buf, 'w', zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    data = zin.read(item.filename)
                    if item.filename == 'xl/worksheets/sheet1.xml':
                        data_str = data.decode('utf-8')
                        # Inject xml:lang on the <t> element of the inline string
                        new_xml = data_str.replace(
                            '<t>Hello</t>', '<t xml:lang="fr-FR">Hello</t>'
                        )
                        data = new_xml.encode('utf-8')
                    zout.writestr(item, data)
        fs = XlsxAnalyzer(out_buf.getvalue(), "lang.xlsx").analyze()
        matches = [f for f in fs.possible_findings
                   if f.criterion_id == "3.1.2" and f.remediation_id == "xlsx_language_of_parts"]
        self.assertGreater(len(matches), 0)


if __name__ == "__main__":
    unittest.main()
