"""Phase I — 6 new strict-deterministic rules.

I1: HTML 1.4.5 Images of Text (inline SVG with <text> children)
I2: HTML 2.5.8 Target Size Minimum (WCAG 2.2 — buttons/links < 24x24 px)
I3: HTML 3.3.8 Accessible Authentication (WCAG 2.2 — password manager blockers)
I4: DOCX 1.3.3 Sensory Characteristics (color/shape/position-only references)
I5: PPTX 1.3.3 Sensory Characteristics (color/shape/position-only references)
I6: XLSX 1.3.3 Sensory Characteristics (color/shape/position-only references)
"""
import io
import unittest
import zipfile

import openpyxl
from pptx import Presentation

from wcag.analyzers.docx_analyzer import DocxAnalyzer
from wcag.analyzers.html_analyzer import HtmlAnalyzer
from wcag.analyzers.pptx_analyzer import PptxAnalyzer
from wcag.analyzers.xlsx_analyzer import XlsxAnalyzer


def _confirmed(fs, criterion: str):
    return [f for f in fs.confirmed_findings if f.criterion_id == criterion]


def _analyze_html(html: str):
    return HtmlAnalyzer(html.encode("utf-8"), "i.html").analyze()


def _build_docx_with_paragraphs(texts):
    """Build a minimal DOCX with the given paragraphs (each a string)."""
    p_xml = "".join(
        f"<w:p><w:r><w:t xml:space='preserve'>{t}</w:t></w:r></w:p>" for t in texts
    )
    doc_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body>{p_xml}</w:body></w:document>"
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
            "</Types>",
        )
        zf.writestr(
            "_rels/.rels",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
            "</Relationships>",
        )
        zf.writestr("word/document.xml", doc_xml)
    return buf.getvalue()


def _build_pptx_with_text_blocks(slides_text):
    """Build a minimal PPTX. slides_text is a list of slides; each slide is a
    list of strings to put in text-box shapes."""
    from pptx.util import Inches
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for texts in slides_text:
        slide = prs.slides.add_slide(blank)
        for i, t in enumerate(texts):
            tb = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(6), Inches(1))
            tb.text_frame.text = t
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_xlsx_with_cells(cell_values):
    """Build a minimal XLSX. cell_values is a list of strings written down col A."""
    wb = openpyxl.Workbook()
    ws = wb.active
    for i, v in enumerate(cell_values, start=1):
        ws.cell(row=i, column=1, value=v)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# I1 — HTML 1.4.5 Inline SVG with <text>
# ─────────────────────────────────────────────────────────────────────────────
class TestHtml145InlineSvgText(unittest.TestCase):
    def test_pass_no_svg(self):
        html = "<!DOCTYPE html><html lang='en'><head><title>x</title></head><body><h1>Hi</h1></body></html>"
        fs = _analyze_html(html)
        self.assertEqual(len(_confirmed(fs, "1.4.5")), 0)

    def test_pass_decorative_svg_no_text(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><svg viewBox='0 0 10 10' aria-hidden='true'><circle cx='5' cy='5' r='4'/></svg></body></html>"""
        fs = _analyze_html(html)
        self.assertEqual(len(_confirmed(fs, "1.4.5")), 0)

    def test_fail_inline_svg_with_text(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><svg width='200' height='40'><text x='10' y='30' font-size='24'>SALE</text></svg></body></html>"""
        fs = _analyze_html(html)
        findings = _confirmed(fs, "1.4.5")
        self.assertEqual(len(findings), 1)
        self.assertIn("svg", findings[0].evidence.lower())

    def test_multiple_svg_text_aggregated(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head><body>
        <svg><text x='0' y='10'>A</text></svg>
        <svg><text x='0' y='10'>B</text></svg>
        </body></html>"""
        fs = _analyze_html(html)
        findings = _confirmed(fs, "1.4.5")
        self.assertEqual(len(findings), 1)
        self.assertIn("2 inline SVG", findings[0].issue)


# ─────────────────────────────────────────────────────────────────────────────
# I2 — HTML 2.5.8 Target Size Minimum (WCAG 2.2)
# ─────────────────────────────────────────────────────────────────────────────
class TestHtml258TargetSize(unittest.TestCase):
    def test_pass_normal_size(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><button style='width:48px;height:48px'>Go</button></body></html>"""
        fs = _analyze_html(html)
        self.assertEqual(len(_confirmed(fs, "2.5.8")), 0)

    def test_fail_inline_style_too_small(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><button style='width:16px;height:16px'>x</button></body></html>"""
        fs = _analyze_html(html)
        findings = _confirmed(fs, "2.5.8")
        self.assertEqual(len(findings), 1)
        self.assertIn("16", findings[0].evidence)

    def test_fail_style_block_too_small(self):
        # Strict rule: only fires when CSS selector references an interactive
        # tag directly (button/a/input/select/textarea). Class-only selectors
        # are out of scope (we can't know if .my-btn is on a <button>).
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title>
        <style>button.tiny { width: 20px; height: 20px; }</style>
        </head><body><button class='tiny'>x</button></body></html>"""
        fs = _analyze_html(html)
        findings = _confirmed(fs, "2.5.8")
        self.assertEqual(len(findings), 1)

    def test_pass_class_only_selector_not_flagged(self):
        # By design: class-only selectors are NOT flagged to avoid false
        # positives (the class might be on a non-interactive icon/element).
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title>
        <style>.tiny { width: 16px; height: 16px; }</style>
        </head><body><button class='tiny'>x</button></body></html>"""
        fs = _analyze_html(html)
        self.assertEqual(len(_confirmed(fs, "2.5.8")), 0)

    def test_pass_only_one_dimension_under_24(self):
        # Rule requires BOTH dimensions < 24 to fire.
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><button style='width:18px;height:30px'>x</button></body></html>"""
        fs = _analyze_html(html)
        self.assertEqual(len(_confirmed(fs, "2.5.8")), 0)


# ─────────────────────────────────────────────────────────────────────────────
# I3 — HTML 3.3.8 Accessible Authentication (WCAG 2.2)
# ─────────────────────────────────────────────────────────────────────────────
class TestHtml338AccessibleAuthentication(unittest.TestCase):
    def test_pass_clean_password_field(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><form><input type='password' name='p' autocomplete='current-password'></form></body></html>"""
        fs = _analyze_html(html)
        self.assertEqual(len(_confirmed(fs, "3.3.8")), 0)

    def test_fail_autocomplete_off(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><form><input type='password' name='p' autocomplete='off'></form></body></html>"""
        fs = _analyze_html(html)
        findings = _confirmed(fs, "3.3.8")
        self.assertEqual(len(findings), 1)
        self.assertIn("autocomplete", findings[0].evidence.lower())

    def test_fail_paste_blocked(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><form><input type='password' name='p' onpaste='return false'></form></body></html>"""
        fs = _analyze_html(html)
        findings = _confirmed(fs, "3.3.8")
        self.assertEqual(len(findings), 1)
        self.assertIn("paste", findings[0].evidence.lower())

    def test_fail_short_maxlength(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><form><input type='password' name='p' maxlength='8'></form></body></html>"""
        fs = _analyze_html(html)
        findings = _confirmed(fs, "3.3.8")
        self.assertEqual(len(findings), 1)
        self.assertIn("maxlength", findings[0].evidence.lower())

    def test_pass_long_maxlength(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><form><input type='password' name='p' maxlength='128' autocomplete='current-password'></form></body></html>"""
        fs = _analyze_html(html)
        self.assertEqual(len(_confirmed(fs, "3.3.8")), 0)


# ─────────────────────────────────────────────────────────────────────────────
# I4 — DOCX 1.3.3 Sensory Characteristics
# ─────────────────────────────────────────────────────────────────────────────
class TestDocx133SensoryCharacteristics(unittest.TestCase):
    def test_pass_clean_text(self):
        data = _build_docx_with_paragraphs([
            "Welcome to the report.",
            "Please review the findings below.",
        ])
        fs = DocxAnalyzer(data, "p.docx").analyze()
        self.assertEqual(len(_confirmed(fs, "1.3.3")), 0)

    def test_fail_color_only(self):
        data = _build_docx_with_paragraphs([
            "Click the red button to continue.",
        ])
        fs = DocxAnalyzer(data, "fail.docx").analyze()
        findings = _confirmed(fs, "1.3.3")
        self.assertEqual(len(findings), 1)
        self.assertEqual(findings[0].remediation_id, "docx_sensory_characteristics")

    def test_fail_position_only(self):
        data = _build_docx_with_paragraphs([
            "See the menu on the right.",
        ])
        fs = DocxAnalyzer(data, "fail.docx").analyze()
        findings = _confirmed(fs, "1.3.3")
        self.assertEqual(len(findings), 1)


# ─────────────────────────────────────────────────────────────────────────────
# I5 — PPTX 1.3.3 Sensory Characteristics
# ─────────────────────────────────────────────────────────────────────────────
class TestPptx133SensoryCharacteristics(unittest.TestCase):
    def test_pass_clean_text(self):
        data = _build_pptx_with_text_blocks([
            ["Welcome", "Agenda for today"],
        ])
        fs = PptxAnalyzer(data, "p.pptx").analyze()
        self.assertEqual(len(_confirmed(fs, "1.3.3")), 0)

    def test_fail_color_only_on_slide(self):
        data = _build_pptx_with_text_blocks([
            ["Click the green icon to advance the slide."],
        ])
        fs = PptxAnalyzer(data, "fail.pptx").analyze()
        findings = _confirmed(fs, "1.3.3")
        self.assertEqual(len(findings), 1)
        self.assertEqual(findings[0].remediation_id, "pptx_sensory_characteristics")


# ─────────────────────────────────────────────────────────────────────────────
# I6 — XLSX 1.3.3 Sensory Characteristics
# ─────────────────────────────────────────────────────────────────────────────
class TestXlsx133SensoryCharacteristics(unittest.TestCase):
    def test_pass_clean_cells(self):
        data = _build_xlsx_with_cells([
            "Region",
            "North",
            "South",
        ])
        fs = XlsxAnalyzer(data, "p.xlsx").analyze()
        self.assertEqual(len(_confirmed(fs, "1.3.3")), 0)

    def test_fail_color_only_in_cell(self):
        data = _build_xlsx_with_cells([
            "Instructions",
            "Click the red button to refresh data.",
        ])
        fs = XlsxAnalyzer(data, "fail.xlsx").analyze()
        findings = _confirmed(fs, "1.3.3")
        self.assertEqual(len(findings), 1)
        self.assertIn("xlsx_sensory_", findings[0].remediation_id)


if __name__ == "__main__":
    unittest.main()
