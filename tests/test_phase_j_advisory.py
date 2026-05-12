"""Phase J — LLM advisory plumbing.

Adds `advisory_payload` field to Finding and populates it for the 4 rules
where deterministic detection finds the issue but only an LLM can draft a
context-aware fix:

J1: HTML 1.1.1 — img missing alt → alt_text advisory
J2: HTML 2.4.4 — generic link text → link_text advisory
J3: DOCX 1.1.1 — image missing alt → alt_text advisory (with paragraph context)
J4: PPTX 1.1.1 — image missing alt → alt_text advisory (with slide context)

Schema verification:
    advisory_payload = {
        "advisory_kind": "alt_text" | "link_text" | "heading_text" | "title_text",
        "target": str,         # stable in-doc reference
        "surface_text": str,   # current bad value (empty when missing)
        "context": str,        # surrounding text (<= 1000 chars)
        "format_hint": "html" | "docx" | "pptx" | "xlsx" | "pdf",
    }
"""
import io
import unittest
import zipfile

from wcag.analyzers.docx_analyzer import DocxAnalyzer
from wcag.analyzers.html_analyzer import HtmlAnalyzer
from wcag.analyzers.pptx_analyzer import PptxAnalyzer


def _all_findings(fs):
    return list(fs.confirmed_findings) + list(fs.possible_findings)


def _findings(fs, criterion: str, advisory_kind: str | None = None):
    out = []
    for f in _all_findings(fs):
        if f.criterion_id != criterion:
            continue
        if advisory_kind is not None:
            if not f.advisory_payload:
                continue
            if f.advisory_payload.get("advisory_kind") != advisory_kind:
                continue
        out.append(f)
    return out


# ─────────────────────────────────────────────────────────────────────────────
# Schema baseline — Finding has the field with default None
# ─────────────────────────────────────────────────────────────────────────────
class TestAdvisoryPayloadFieldExists(unittest.TestCase):
    def test_field_default_is_none_for_findings_without_advisory(self):
        # A page with NO images and NO links → 1.1.1 / 2.4.4 will not fire,
        # but other rules will. Verify advisory_payload defaults to None.
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><h1>Hi</h1><p>Plain text only.</p></body></html>"""
        fs = HtmlAnalyzer(html.encode("utf-8"), "x.html").analyze()
        for f in _all_findings(fs):
            # Findings that don't opt in must have None.
            if f.criterion_id not in ("1.1.1", "2.4.4"):
                self.assertIsNone(
                    f.advisory_payload,
                    f"finding {f.criterion_id}/{f.remediation_id} should not "
                    f"have advisory_payload but does: {f.advisory_payload}"
                )

    def test_to_dict_serializes_advisory_payload_key(self):
        # Even when None, the key should be present in the dict.
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><img src='x.png'><a href='/y'>Click here</a></body></html>"""
        fs = HtmlAnalyzer(html.encode("utf-8"), "x.html").analyze()
        for f in _all_findings(fs):
            d = f.to_dict()
            self.assertIn("advisory_payload", d)


# ─────────────────────────────────────────────────────────────────────────────
# J1 — HTML 1.1.1 missing alt
# ─────────────────────────────────────────────────────────────────────────────
class TestJ1HtmlImageAltAdvisory(unittest.TestCase):
    def test_advisory_payload_populated_with_context(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body>
          <h2>Quarterly Sales Performance</h2>
          <p>Revenue grew 23% in Q3 driven by APAC expansion.</p>
          <img src='https://cdn.example.com/q3-chart.png'>
        </body></html>"""
        fs = HtmlAnalyzer(html.encode("utf-8"), "x.html").analyze()
        findings = _findings(fs, "1.1.1", "alt_text")
        self.assertEqual(len(findings), 1)
        ap = findings[0].advisory_payload
        self.assertEqual(ap["advisory_kind"], "alt_text")
        self.assertEqual(ap["format_hint"], "html")
        self.assertEqual(ap["surface_text"], "")
        # target should be the src URL
        self.assertEqual(ap["target"], "https://cdn.example.com/q3-chart.png")
        # context should contain the heading and paragraph text
        self.assertIn("Quarterly Sales Performance", ap["context"])
        self.assertIn("Revenue grew 23%", ap["context"])

    def test_decorative_image_has_no_finding_and_no_advisory(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><img src='x.png' alt='' role='presentation'></body></html>"""
        fs = HtmlAnalyzer(html.encode("utf-8"), "x.html").analyze()
        self.assertEqual(len(_findings(fs, "1.1.1", "alt_text")), 0)


# ─────────────────────────────────────────────────────────────────────────────
# J2 — HTML 2.4.4 generic link text
# ─────────────────────────────────────────────────────────────────────────────
class TestJ2HtmlLinkTextAdvisory(unittest.TestCase):
    def test_generic_link_gets_advisory_with_context_and_href(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body>
          <p>To download the latest Q3 financial report, please
          <a href='https://example.com/reports/q3-2026.pdf'>click here</a>.</p>
        </body></html>"""
        fs = HtmlAnalyzer(html.encode("utf-8"), "x.html").analyze()
        findings = _findings(fs, "2.4.4", "link_text")
        self.assertEqual(len(findings), 1)
        ap = findings[0].advisory_payload
        self.assertEqual(ap["advisory_kind"], "link_text")
        self.assertEqual(ap["format_hint"], "html")
        self.assertEqual(ap["surface_text"], "click here")
        self.assertEqual(ap["target"], "https://example.com/reports/q3-2026.pdf")
        self.assertIn("Q3 financial report", ap["context"])

    def test_descriptive_link_no_advisory(self):
        html = """<!DOCTYPE html><html lang='en'><head><title>x</title></head>
        <body><a href='/docs'>Read the deployment handbook</a></body></html>"""
        fs = HtmlAnalyzer(html.encode("utf-8"), "x.html").analyze()
        self.assertEqual(len(_findings(fs, "2.4.4", "link_text")), 0)


# ─────────────────────────────────────────────────────────────────────────────
# J3 — DOCX 1.1.1 missing alt
# ─────────────────────────────────────────────────────────────────────────────
def _docx_with_image_xml(extra_pre_paragraphs: str = "", descr_attr: str = ""):
    """Build a minimal DOCX with one inline image (and optional surrounding
    paragraphs for context). `descr_attr` like ' descr="X"' or '' for missing."""
    document_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
            xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
            xmlns:pic="http://schemas.openxmlformats.org/drawingml/2006/picture">
  <w:body>
    {extra_pre_paragraphs}
    <w:p>
      <w:r>
        <w:drawing>
          <wp:inline>
            <wp:extent cx="100" cy="100"/>
            <wp:docPr id="1" name="Picture 1"{descr_attr}/>
            <wp:cNvGraphicFramePr/>
            <a:graphic>
              <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/picture">
                <pic:pic>
                  <pic:nvPicPr>
                    <pic:cNvPr id="1" name="Picture 1"/>
                    <pic:cNvPicPr/>
                  </pic:nvPicPr>
                  <pic:blipFill><a:blip r:embed="rId1"/></pic:blipFill>
                  <pic:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="100" cy="100"/></a:xfrm>
                    <a:prstGeom prst="rect"/></pic:spPr>
                </pic:pic>
              </a:graphicData>
            </a:graphic>
          </wp:inline>
        </w:drawing>
      </w:r>
    </w:p>
    <w:sectPr><w:pgSz w:w="12240" w:h="15840"/>
      <w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440"
               w:header="720" w:footer="720" w:gutter="0"/></w:sectPr>
  </w:body>
</w:document>"""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml",
                    """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""")
        zf.writestr("_rels/.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""")
        zf.writestr("word/document.xml", document_xml)
    return buffer.getvalue()


class TestJ3DocxImageAltAdvisory(unittest.TestCase):
    def test_image_missing_alt_carries_advisory_with_paragraph_context(self):
        pre = """<w:p><w:pPr><w:pStyle w:val="Heading1"/></w:pPr>
                  <w:r><w:t>Q3 Financial Highlights</w:t></w:r></w:p>
                <w:p><w:r><w:t>Revenue grew 23% in Q3.</w:t></w:r></w:p>"""
        data = _docx_with_image_xml(extra_pre_paragraphs=pre)
        fs = DocxAnalyzer(data, "t.docx").analyze()
        findings = _findings(fs, "1.1.1", "alt_text")
        self.assertEqual(len(findings), 1)
        ap = findings[0].advisory_payload
        self.assertEqual(ap["advisory_kind"], "alt_text")
        self.assertEqual(ap["format_hint"], "docx")
        self.assertEqual(ap["surface_text"], "")
        self.assertIn("Q3 Financial Highlights", ap["context"])
        self.assertIn("Revenue grew 23%", ap["context"])

    def test_image_with_alt_text_has_no_advisory_finding(self):
        data = _docx_with_image_xml(descr_attr=' descr="A bar chart of Q3 revenue"')
        fs = DocxAnalyzer(data, "t.docx").analyze()
        # No 1.1.1 advisory should fire because alt text is present
        self.assertEqual(len(_findings(fs, "1.1.1", "alt_text")), 0)


# ─────────────────────────────────────────────────────────────────────────────
# J4 — PPTX 1.1.1 missing alt
# ─────────────────────────────────────────────────────────────────────────────
def _pptx_with_picture(descr: str | None = None):
    """Build a PPTX with a slide containing a title and a picture.

    descr=None  → simulate missing alt text (delete cNvPr@descr)
    descr=""    → empty descr (explicit empty)
    descr="..." → set descr to that value
    """
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image
    import tempfile
    import os

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Q3 Financial Highlights"
    img_path = tempfile.mktemp(suffix=".png")
    Image.new("RGB", (50, 50), "blue").save(img_path)
    try:
        pic = slide.shapes.add_picture(img_path, Inches(1), Inches(1), Inches(1), Inches(1))
        cNvPr = pic._element.nvPicPr.cNvPr
        if descr is None:
            # python-pptx auto-fills descr with the source filename — clear it.
            if "descr" in cNvPr.attrib:
                del cNvPr.attrib["descr"]
        else:
            cNvPr.set("descr", descr)
    finally:
        try:
            os.remove(img_path)
        except OSError:
            pass
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestJ4PptxImageAltAdvisory(unittest.TestCase):
    def test_pptx_image_missing_alt_carries_advisory_with_slide_context(self):
        data = _pptx_with_picture(descr=None)
        fs = PptxAnalyzer(data, "p.pptx").analyze()
        findings = _findings(fs, "1.1.1", "alt_text")
        self.assertEqual(len(findings), 1)
        ap = findings[0].advisory_payload
        self.assertEqual(ap["advisory_kind"], "alt_text")
        self.assertEqual(ap["format_hint"], "pptx")
        self.assertEqual(ap["surface_text"], "")
        # Context should mention slide title
        self.assertIn("Q3 Financial Highlights", ap["context"])
        # Target should reference slide + shape id
        self.assertTrue(ap["target"].startswith("slide1/shape"))

    def test_pptx_image_with_alt_text_has_no_advisory(self):
        data = _pptx_with_picture(descr="A bar chart showing Q3 revenue")
        fs = PptxAnalyzer(data, "p.pptx").analyze()
        self.assertEqual(len(_findings(fs, "1.1.1", "alt_text")), 0)


if __name__ == "__main__":
    unittest.main()
